#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for Umuganda Social Impact Tracking Platform
Using python-pptx library
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define Umuganda color scheme
ONYX = RGBColor(13, 10, 11)  # #0d0a0bff
CHARCOAL_BLUE = RGBColor(69, 73, 85)  # #454955ff
LAVENDER_MIST = RGBColor(243, 239, 245)  # #f3eff5ff
LIME_MOSS = RGBColor(114, 176, 29)  # #72b01dff
GREEN = RGBColor(63, 125, 32)  # #3f7d20ff
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ONYX
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = LIME_MOSS
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LAVENDER_MIST
    p.alignment = PP_ALIGN.CENTER
    
    # Footer date
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = f"December 2025"
    p.font.size = Pt(12)
    p.font.color.rgb = CHARCOAL_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_items, include_image=False):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LAVENDER_MIST
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = LIME_MOSS
    title_shape.line.color.rgb = LIME_MOSS
    
    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(10)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = ONYX
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_two_column_slide(title, left_content, right_content):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LAVENDER_MIST
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = LIME_MOSS
    title_shape.line.color.rgb = LIME_MOSS
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(6)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(6))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = ONYX
        p.space_before = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(6))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = ONYX
        p.space_before = Pt(6)

# Slide 1: Title Slide
add_title_slide(
    "Umuganda",
    "Social Impact Tracking Platform for Rwanda"
)

# Slide 2: Problem Statement
add_content_slide(
    "Problem Statement",
    [
        "• Manual tracking: Community service activities recorded manually without centralization",
        "• Lack of transparency: No visibility into impact metrics or community contributions",
        "• Low engagement: Communities lack incentive to report and celebrate contributions",
        "• No analytics: Absence of data-driven insights for decision-making at district level",
        "• Siloed information: Each cell operates independently without sharing knowledge",
        "• Accountability gaps: Difficult to verify and measure social impact"
    ]
)

# Slide 3: Solution Overview
add_content_slide(
    "Solution Overview",
    [
        "📱 Unified Platform: Web-based application for social media-style activity sharing",
        "🎯 Impact Measurement: Real-time analytics dashboard tracking community contributions",
        "👥 Hierarchical Access: Cell members, moderators, and district viewers with role-based permissions",
        "📊 Rich Engagement: Comments, reactions, reposts to foster community participation",
        "🏆 Recognition System: Trending activities, popular posts, engagement leaderboards",
        "🔍 Data Insights: Comprehensive analytics showing impact by location, category, and time"
    ]
)

# Slide 4: Key Features
add_two_column_slide(
    "Core Features",
    [
        "User Management:",
        "• Registration & authentication",
        "• Role-based access control",
        "• Profile management",
        "",
        "Activity Sharing:",
        "• Text posts with hashtags",
        "• Image uploads (up to 5)",
        "• Location tagging",
        "• Category classification"
    ],
    [
        "Engagement Tools:",
        "• Comments on posts",
        "• 4 reaction types",
        "• Repost/share capability",
        "• Trending detection",
        "",
        "Analytics & Insights:",
        "• National statistics",
        "• Category breakdown",
        "• Popular hashtags",
        "• Community leaderboards"
    ]
)

# Slide 5: System Architecture - Use Case Diagram
add_content_slide(
    "Use Case Diagram",
    [
        "👤 Actors:",
        "   • Community Member: Create posts, comment, react, repost",
        "   • Cell Moderator: Manage posts, moderate comments, oversee cell activities",
        "   • District Viewer: View analytics, generate reports, access all district data",
        "",
        "🎯 Primary Use Cases:",
        "   1. Register & Login → Access control",
        "   2. Create Post → Share activity (text/images)",
        "   3. Engage Post → Comment, react, repost",
        "   4. View Feed → Discover activities with filters",
        "   5. View Analytics → Track impact metrics"
    ]
)

# Slide 6: Data Flow - Sequence Diagram
add_content_slide(
    "Post Creation Workflow",
    [
        "📝 High-Level Sequence:",
        "",
        "1. User fills compose form (text/images, location, category, hashtags)",
        "2. Frontend validates input (length, image count, etc.)",
        "3. API receives request + validates with Zod schemas",
        "4. Factory pattern creates appropriate post type (Text/Image)",
        "5. Builder pattern constructs post with all metadata",
        "6. Prisma ORM saves to PostgreSQL database",
        "7. Analytics engine updates metrics",
        "8. Response sent with post object",
        "9. Frontend updates UI and shows confirmation",
        "",
        "⚡ Performance: <200ms database operation + validation"
    ]
)

# Save presentation
output_path = "Umuganda_Platform_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"📁 Location: {output_path}")
